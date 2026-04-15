"""Chip Day 2026 slide upload server.

Runs behind nginx reverse proxy at /2026/upload/.
Speakers upload slides to their designated slot; files are stored per-slot
and served from the protected /2026/ area. After each upload, the session
deck is automatically re-merged.
"""

import os
import json
import copy
import logging
from pathlib import Path
from threading import Thread
from flask import Flask, request, redirect, jsonify
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

app = Flask(__name__)
logging.basicConfig(level=logging.INFO)
log = logging.getLogger("chipday-upload")

UPLOAD_DIR = Path("/var/www/chipday.dk/2026/uploads")
MERGED_DIR = Path("/var/www/chipday.dk/2026/merged")
PRES_DIR = Path("/var/www/chipday.dk/2026/presentations")
LOGO_DIR = Path("/var/www/chipday.dk/assets/logos")
QUIZ_FILE = UPLOAD_DIR / "quiz-status.json"
MAX_FILE_SIZE = 200 * 1024 * 1024  # 200 MB
ALLOWED_EXTENSIONS = {".pptx", ".ppt", ".pdf", ".key"}
LOGO_EXTENSIONS = {".png", ".svg", ".jpg", ".jpeg", ".webp"}

# Relationship types to skip when cloning slides (already in destination)
SKIP_RELTYPES = {RT.SLIDE_LAYOUT, RT.SLIDE_MASTER, RT.THEME}

SESSIONS = {
    "session-1": {
        "name": "Chip Education in Denmark",
        "slots": ["s1-intro", "s1-dam", "s1-lynggaard", "s1-schoeberl", "s1-moradi", "s1-rolver"],
    },
    "session-2": {
        "name": "Technical Presentations I",
        "slots": ["s2-nour", "s2-dahl", "s2-marquart"],
    },
    "session-3": {
        "name": "Technical Presentations II",
        "slots": ["s3-smedegaard", "s3-saerkjaer", "s3-lind"],
    },
}

SLOTS = {
    # Session 1: Chip Education in Denmark
    "s1-intro":       "Welcome & Introduction",
    "s1-dam":         "Keynote: The Importance of IC Development",
    "s1-lynggaard":   "Analogue Courses",
    "s1-schoeberl":   "Digital Courses + First Student Tape-Out",
    "s1-moradi":      "Research and Chip Courses at SDU",
    "s1-rolver":      "The Industry Master at DTU",
    # Session 2: Technical Presentations I
    "s2-nour":        "Lotus Microsystems",
    "s2-dahl":        "EDA Startup",
    "s2-marquart":    "DSP Audio Processing in FPGA for Jabra Product",
    # Session 3: Technical Presentations II
    "s3-smedegaard":  "Keynote: Opportunities for DK Chip Design",
    "s3-saerkjaer":   "Quantum Foundry",
    "s3-lind":        "Ultra High Speed Communication",
}


def slot_dir(slot_id):
    return UPLOAD_DIR / slot_id


def get_current_file(slot_id):
    d = slot_dir(slot_id)
    if not d.exists():
        return None
    files = [f for f in d.iterdir() if f.is_file()]
    return files[0] if files else None


def get_pptx_file(slot_id):
    """Get the .pptx file for a slot, or None."""
    f = get_current_file(slot_id)
    if f and f.suffix.lower() == ".pptx":
        return f
    return None


def merge_session(session_id):
    """Merge all .pptx files for a session into one deck."""
    session = SESSIONS.get(session_id)
    if not session:
        return

    pptx_files = []
    for slot_id in session["slots"]:
        f = get_pptx_file(slot_id)
        if f:
            pptx_files.append(f)

    if not pptx_files:
        log.info(f"No .pptx files for {session_id}, skipping merge")
        return

    MERGED_DIR.mkdir(parents=True, exist_ok=True)
    output = MERGED_DIR / f"{session_id}.pptx"

    log.info(f"Merging {session_id}: {[f.name for f in pptx_files]}")

    dst = Presentation(str(pptx_files[0]))

    for pptx_path in pptx_files[1:]:
        src = Presentation(str(pptx_path))
        for src_slide in src.slides:
            # Find best matching layout
            src_layout_name = src_slide.slide_layout.name
            dst_layout = dst.slide_layouts[0]
            for layout in dst.slide_layouts:
                if layout.name == src_layout_name:
                    dst_layout = layout
                    break

            dst_slide = dst.slides.add_slide(dst_layout)

            # Remove default placeholders
            for ph in list(dst_slide.placeholders):
                dst_slide.shapes._spTree.remove(ph._element)

            # Copy slide content
            src_cSld = src_slide._element.find(
                '{http://schemas.openxmlformats.org/presentationml/2006/main}cSld'
            )
            dst_cSld = dst_slide._element.find(
                '{http://schemas.openxmlformats.org/presentationml/2006/main}cSld'
            )
            if src_cSld is not None and dst_cSld is not None:
                dst_slide._element.replace(dst_cSld, copy.deepcopy(src_cSld))

            # Copy media relationships
            for rel in src_slide.part.rels.values():
                if rel.reltype in SKIP_RELTYPES:
                    continue
                if rel.is_external:
                    dst_slide.part.rels.get_or_add_ext_rel(
                        rel.reltype, rel.target_ref
                    )
                else:
                    try:
                        dst_slide.part.rels.get_or_add(
                            rel.reltype, rel.target_part
                        )
                    except Exception:
                        pass

    dst.save(str(output))
    log.info(f"Saved {output} ({len(dst.slides)} slides)")


def find_session_for_slot(slot_id):
    """Return the session_id that contains this slot."""
    for session_id, session in SESSIONS.items():
        if slot_id in session["slots"]:
            return session_id
    return None


def trigger_merge(slot_id):
    """Trigger a background merge for the session containing this slot."""
    session_id = find_session_for_slot(slot_id)
    if session_id:
        thread = Thread(target=merge_session, args=(session_id,))
        thread.start()


@app.route("/2026/upload/status")
def status():
    """Return JSON with upload status for all slots and merged decks."""
    result = {"slots": {}, "merged": {}}
    for slot_id, title in SLOTS.items():
        f = get_current_file(slot_id)
        result["slots"][slot_id] = {
            "title": title,
            "uploaded": f is not None,
            "filename": f.name if f else None,
        }
    for session_id, session in SESSIONS.items():
        merged = MERGED_DIR / f"{session_id}.pptx"
        result["merged"][session_id] = {
            "name": session["name"],
            "available": merged.exists(),
        }
    return jsonify(result)


@app.route("/2026/upload/<slot_id>", methods=["POST"])
def upload(slot_id):
    if slot_id not in SLOTS:
        return "Unknown slot", 404

    if "file" not in request.files:
        return "No file provided", 400

    f = request.files["file"]
    if not f.filename:
        return "No file selected", 400

    ext = os.path.splitext(f.filename)[1].lower()
    if ext not in ALLOWED_EXTENSIONS:
        return f"File type {ext} not allowed. Use: {', '.join(ALLOWED_EXTENSIONS)}", 400

    # Clear previous uploads for this slot
    d = slot_dir(slot_id)
    d.mkdir(parents=True, exist_ok=True)
    for old in d.iterdir():
        old.unlink()

    # Save with original filename
    safe_name = f.filename.replace("/", "_").replace("\\", "_")
    f.save(d / safe_name)
    log.info(f"Upload: {slot_id} <- {safe_name}")

    # Trigger merge in background
    trigger_merge(slot_id)

    return redirect("/2026/#uploads")


def load_quiz_status():
    if QUIZ_FILE.exists():
        return json.loads(QUIZ_FILE.read_text())
    return {}


def save_quiz_status(data):
    QUIZ_FILE.parent.mkdir(parents=True, exist_ok=True)
    QUIZ_FILE.write_text(json.dumps(data))


@app.route("/2026/upload/quiz-status")
def quiz_status():
    return jsonify(load_quiz_status())


@app.route("/2026/upload/quiz-status/<slot_id>", methods=["POST"])
def quiz_toggle(slot_id):
    data = load_quiz_status()
    body = request.get_json(silent=True) or {}
    data[slot_id] = bool(body.get("checked", False))
    save_quiz_status(data)
    return jsonify({"ok": True})


@app.route("/api/logo-status")
@app.route("/logos/upload/status")
def logo_status():
    """Return JSON with logo upload status for all companies."""
    result = {}
    if LOGO_DIR.exists():
        for f in LOGO_DIR.iterdir():
            if f.is_file() and f.suffix.lower() in LOGO_EXTENSIONS:
                slug = f.stem
                result[slug] = {"uploaded": True, "filename": f.name}
    return jsonify(result)


@app.route("/logos/upload/<slug>", methods=["POST"])
def logo_upload(slug):
    if "file" not in request.files:
        return "No file provided", 400

    f = request.files["file"]
    if not f.filename:
        return "No file selected", 400

    ext = os.path.splitext(f.filename)[1].lower()
    if ext not in LOGO_EXTENSIONS:
        return f"File type {ext} not allowed. Use: {', '.join(LOGO_EXTENSIONS)}", 400

    LOGO_DIR.mkdir(parents=True, exist_ok=True)

    # Remove any existing logo for this slug
    for old in LOGO_DIR.iterdir():
        if old.is_file() and old.stem == slug:
            old.unlink()

    # Save as slug + extension (e.g. ic-works.png)
    dest = LOGO_DIR / f"{slug}{ext}"
    f.save(dest)
    log.info(f"Logo upload: {slug} <- {dest.name}")

    return redirect("/logos/")


# Mapping from slot id to clean filename for the presentations folder
PRES_NAMES = {
    "s1-intro":      "session-1_education/s1-00-Introduction",
    "s1-dam":        "session-1_education/s1-01-Brian_Dam_Pedersen-Keynote_IC_Development",
    "s1-lynggaard":  "session-1_education/s1-02-Per_Lynggaard-Analogue_Courses",
    "s1-schoeberl":  "session-1_education/s1-03-Martin_Schoeberl-First_Student_Tape_Out",
    "s1-moradi":     "session-1_education/s1-04-Farshad_Moradi-Chip_Courses_SDU",
    "s1-rolver":     "session-1_education/s1-05-Katrine_Rolver-Industry_Master_DTU",
    "s2-nour":       "session-2_technical_I/s2-01-Martin_Lantz-Lotus_Microsystems",
    "s2-dahl":       "session-2_technical_I/s2-02-Nicolai_Dahl-IC_Optimize_EDA",
    "s2-marquart":   "session-2_technical_I/s2-03-Jimmi_Marquart-DSP_Audio_FPGA",
    "s3-smedegaard": "session-3_technical_II/s3-01-Michael_Smedegaard-Opportunities_DK_Chip",
    "s3-saerkjaer":  "session-3_technical_II/s3-02-Tobias_Saerkjaer-Quantum_Foundry",
    "s3-lind":       "session-3_technical_II/s3-03-Rasmus_Lind-Ultra_High_Speed_Comms",
}

PRES_SUBDIRS = [
    "session-1_education",
    "session-2_technical_I",
    "session-3_technical_II",
]


def rebuild_presentations():
    """Rebuild the organized presentations folder and zip from uploads."""
    import shutil
    import zipfile

    # Clear and recreate
    if PRES_DIR.exists():
        shutil.rmtree(PRES_DIR)
    for subdir in PRES_SUBDIRS:
        (PRES_DIR / subdir).mkdir(parents=True, exist_ok=True)

    copied = []
    for slot_id, dest_path in PRES_NAMES.items():
        src = get_current_file(slot_id)
        if not src:
            continue
        ext = src.suffix
        dst = PRES_DIR / f"{dest_path}{ext}"
        shutil.copy2(str(src), str(dst))
        copied.append(dst.name)

    # Build zip
    zip_path = PRES_DIR.parent / "presentations.zip"
    with zipfile.ZipFile(str(zip_path), "w", zipfile.ZIP_DEFLATED) as zf:
        for root, dirs, files in os.walk(str(PRES_DIR)):
            for fname in sorted(files):
                fpath = Path(root) / fname
                arcname = f"presentations/{fpath.relative_to(PRES_DIR)}"
                zf.write(str(fpath), arcname)

    log.info(f"Rebuilt presentations: {len(copied)} files, zip={zip_path}")
    return copied


@app.route("/2026/upload/rebuild-presentations", methods=["POST"])
def rebuild_presentations_endpoint():
    """Rebuild the organized presentations folder from current uploads."""
    try:
        copied = rebuild_presentations()
        return jsonify({"ok": True, "files": len(copied), "names": copied})
    except Exception as e:
        log.exception("Rebuild failed")
        return jsonify({"ok": False, "error": str(e)}), 500


if __name__ == "__main__":
    app.run(host="127.0.0.1", port=5026, debug=True)
